# -*- coding: utf-8 -*-
"""
PPT处理器 - 集成PPT/PPTX解析和内容提取功能
"""

import os
import json
import time
import pickle
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from uuid import uuid4
from loguru import logger

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx库未安装，PPTX文件处理功能将不可用")

try:
    import win32com.client
    import pythoncom
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False
    logger.warning("pywin32库未安装，PPT文件处理功能将不可用")

from .config import OUTPUT_DIR, PICKLES_DIR, PROMPTS
from .ocr_engine import OCREngine
try:
    from .llm_processor import LLMProcessor
except Exception:
    LLMProcessor = None


class PPTProcessor:
    """PPT处理器，负责PPT/PPTX解析和内容提取"""
    
    def __init__(self, ocr_engine: Optional[OCREngine] = None):
        """
        初始化PPT处理器
        """
        self.llm_processor = LLMProcessor() if LLMProcessor else None
        self.ocr_engine = ocr_engine
        self.supported_formats = []
        
        if PPTX_AVAILABLE:
            self.supported_formats.append('.pptx')
        if PPT_AVAILABLE:
            self.supported_formats.append('.ppt')
        
        logger.info(f"PPT处理器初始化完成，支持格式: {self.supported_formats}")
        
    def can_process(self, file_path: str) -> bool:
        """
        检查是否可以处理该文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            是否可以处理
        """
        file_ext = Path(file_path).suffix.lower()
        return file_ext in self.supported_formats
    
    def process_ppt(self, ppt_path: str, output_name: str = None) -> Dict:
        """
        处理PPT文件
        
        Args:
            ppt_path: PPT文件路径
            output_name: 输出名称，如果为None则使用文件名
            
        Returns:
            处理结果字典
        """
        if not Path(ppt_path).exists():
            return {'status': 'error', 'message': f'PPT文件不存在: {ppt_path}'}
        
        if output_name is None:
            output_name = Path(ppt_path).stem
        
        if not self.can_process(ppt_path):
            return {'status': 'error', 'message': f'不支持的文件格式: {Path(ppt_path).suffix}'}
        
        try:
            # 创建输出目录
            output_path = OUTPUT_DIR / output_name
            output_path.mkdir(exist_ok=True)
            
            # 处理PPT
            result = self._process_ppt_content(ppt_path, output_path, output_name)
            
            # 清理临时文件
            self._cleanup_temp_files(output_path)
            
            return result
            
        except Exception as e:
            logger.error(f"PPT处理失败: {e}")
            return {'status': 'error', 'message': str(e)}
    
    def _process_ppt_content(self, ppt_path: str, output_path: Path, output_name: str) -> Dict:
        """处理PPT内容"""
        try:
            file_ext = Path(ppt_path).suffix.lower()
            
            if file_ext == '.pptx':
                return self._process_pptx_file(ppt_path, output_path, output_name)
            elif file_ext == '.ppt':
                return self._process_ppt_file(ppt_path, output_path, output_name)
            else:
                return {'status': 'error', 'message': f'不支持的文件格式: {file_ext}'}
                
        except Exception as e:
            logger.error(f"PPT内容处理失败: {e}")
            return {'status': 'error', 'message': str(e)}
    
    def _process_pptx_file(self, pptx_path: str, output_path: Path, output_name: str) -> Dict:
        """处理PPTX文件"""
        try:
            presentation = Presentation(pptx_path)
            total_slides = len(presentation.slides)
            
            logger.info(f"开始处理PPTX: {pptx_path}, 共{total_slides}页")
            
            all_texts = []
            slide_texts = []
            
            for slide_num, slide in enumerate(presentation.slides):
                logger.info(f"处理第{slide_num + 1}页")
                
                slide_text = self._extract_slide_text(slide, slide_num)
                slide_texts.append(slide_text)
                all_texts.extend(slide_text)

                # 若该页文本过少，尝试对图片进行OCR补充
                if self.ocr_engine and len(''.join(slide_text).strip()) < 10:
                    ocr_extra = self._ocr_pptx_slide_images(slide, output_path, slide_num)
                    if ocr_extra:
                        slide_texts[-1].extend(ocr_extra)
                        all_texts.extend(ocr_extra)
                
                # 进度更新
                progress = (slide_num + 1) / total_slides * 100
                logger.info(f"处理进度: {progress:.1f}%")
            
            # 服务器端不执行LLM，返回空占位
            summary_result = {
                'summary': '',
                'keywords': [],
                'hybrid_summary': '',
                'markdown_content': '',
                'part_summaries': []
            }
            
            # 保存结果
            result = {
                'status': 'success',
                'file_path': pptx_path,
                'output_name': output_name,
                'total_slides': total_slides,
                'total_text_length': len(''.join(all_texts)),
                'slide_texts': slide_texts,
                'summary': summary_result.get('summary', ''),
                'keywords': summary_result.get('keywords', []),
                'processing_time': time.time(),
                'file_format': 'pptx'
            }
            
            # 保存到pickle文件
            self._save_result(result, output_path, output_name)
            
            logger.info(f"PPTX处理完成: {output_name}")
            return result
            
        except Exception as e:
            logger.error(f"PPTX处理失败: {e}")
            return {'status': 'error', 'message': str(e)}
    
    def _process_ppt_file(self, ppt_path: str, output_path: Path, output_name: str) -> Dict:
        """处理PPT文件（使用COM接口）"""
        try:
            # 初始化COM
            pythoncom.CoInitialize()
            
            # 创建PowerPoint应用实例
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False
            
            try:
                # 打开演示文稿
                presentation = powerpoint.Presentations.Open(ppt_path)
                total_slides = presentation.Slides.Count
                
                logger.info(f"开始处理PPT: {ppt_path}, 共{total_slides}页")
                
                all_texts = []
                slide_texts = []
                
                for slide_num in range(1, total_slides + 1):
                    slide = presentation.Slides(slide_num)
                    logger.info(f"处理第{slide_num}页")
                    
                    slide_text = self._extract_ppt_slide_text(slide, slide_num)
                    slide_texts.append(slide_text)
                    all_texts.extend(slide_text)
                    
                    # 进度更新
                    progress = slide_num / total_slides * 100
                    logger.info(f"处理进度: {progress:.1f}%")
                
                # 关闭演示文稿
                presentation.Close()
                
                # 服务器端不执行LLM，返回空占位
                summary_result = {
                    'summary': '',
                    'keywords': [],
                    'hybrid_summary': '',
                    'markdown_content': '',
                    'part_summaries': []
                }
                
                # 保存结果
                result = {
                    'status': 'success',
                    'file_path': ppt_path,
                    'output_name': output_name,
                    'total_slides': total_slides,
                    'total_text_length': len(''.join(all_texts)),
                    'slide_texts': slide_texts,
                    'summary': summary_result.get('summary', ''),
                    'keywords': summary_result.get('keywords', []),
                    'hybrid_summary': summary_result.get('hybrid_summary', ''),
                    'markdown_content': summary_result.get('markdown_content', ''),
                    'part_summaries': summary_result.get('part_summaries', []),
                    'processing_time': time.time(),
                    'file_format': 'ppt'
                }
                
                # 保存到pickle文件
                self._save_result(result, output_path, output_name)
                
                logger.info(f"PPT处理完成: {output_name}")
                return result
                
            finally:
                # 关闭PowerPoint应用
                powerpoint.Quit()
                pythoncom.CoUninitialize()
                
        except Exception as e:
            logger.error(f"PPT处理失败: {e}")
            return {'status': 'error', 'message': str(e)}
    
    def _extract_slide_text(self, slide, slide_num: int) -> List[str]:
        """从PPTX幻灯片提取文本"""
        texts = []
        
        try:
            # 提取形状中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
                        logger.debug(f"幻灯片{slide_num + 1} - 形状文本: {text[:50]}...")
                
                # 提取表格中的文本
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                texts.append(cell.text.strip())
                                logger.debug(f"幻灯片{slide_num + 1} - 表格文本: {cell.text[:50]}...")
                
                # 提取文本框中的文本
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text.strip())
                            logger.debug(f"幻灯片{slide_num + 1} - 段落文本: {paragraph.text[:50]}...")
            
            logger.info(f"幻灯片{slide_num + 1}提取到{len(texts)}个文本片段")
            return texts
            
        except Exception as e:
            logger.error(f"提取幻灯片{slide_num + 1}文本失败: {e}")
            return []

    def _ocr_pptx_slide_images(self, slide, output_path: Path, slide_num: int) -> List[str]:
        """对幻灯片中的图片形状进行OCR，返回补充文本"""
        texts: List[str] = []
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in getattr(slide, 'shapes', []):
                try:
                    if getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, 'image'):
                        image_blob = shape.image.blob
                        if not image_blob:
                            continue
                        img_name = f"slide_{slide_num + 1}_img_{shape.shape_id}.png"
                        img_path = output_path / img_name
                        with open(img_path, 'wb') as f:
                            f.write(image_blob)
                        if self.ocr_engine:
                            ocr_items = self.ocr_engine.extract_text_direct(str(img_path))
                            if ocr_items:
                                merged = "\n".join([t.get('text', '') for t in ocr_items if t.get('text')])
                                if merged.strip():
                                    texts.append(merged)
                except Exception as _e:
                    logger.debug(f"OCR图片失败(第{slide_num + 1}页某图): {_e}")
            if texts:
                logger.info(f"幻灯片{slide_num + 1}通过图片OCR补充到{len(texts)}段文本")
            return texts
        except Exception as e:
            logger.warning(f"OCR图片补充失败(第{slide_num + 1}页): {e}")
            return []
    
    def _extract_ppt_slide_text(self, slide, slide_num: int) -> List[str]:
        """从PPT幻灯片提取文本（COM接口）"""
        texts = []
        
        try:
            # 提取形状中的文本
            for shape in slide.Shapes:
                if hasattr(shape, "TextFrame") and shape.TextFrame:
                    text = shape.TextFrame.TextRange.Text.strip()
                    if text:
                        texts.append(text)
                        logger.debug(f"幻灯片{slide_num} - 形状文本: {text[:50]}...")
                
                # 提取表格中的文本
                if shape.HasTable:
                    table = shape.Table
                    for row in range(1, table.Rows.Count + 1):
                        for col in range(1, table.Columns.Count + 1):
                            cell = table.Cell(row, col)
                            if cell.Shape.TextFrame:
                                text = cell.Shape.TextFrame.TextRange.Text.strip()
                                if text:
                                    texts.append(text)
                                    logger.debug(f"幻灯片{slide_num} - 表格文本: {text[:50]}...")
            
            logger.info(f"幻灯片{slide_num}提取到{len(texts)}个文本片段")
            return texts
            
        except Exception as e:
            logger.error(f"提取幻灯片{slide_num}文本失败: {e}")
            return []
    
    def _generate_summary(self, texts: List[str]) -> Dict:
        """生成摘要和关键词"""
        try:
            if not texts:
                return {
                    'summary': '',
                    'keywords': [],
                    'hybrid_summary': '',
                    'markdown_content': '',
                    'part_summaries': []
                }
            
            # 合并所有文本
            combined_text = '\n'.join(texts)
            
            # 使用LLM生成各种类型的内容
            summary = self.llm_processor.generate_summary(combined_text)
            keywords = self.llm_processor.extract_keywords(combined_text)
            hybrid_summary = self.llm_processor.generate_hybrid_summary(combined_text)
            markdown_content = self.llm_processor.convert_to_markdown(combined_text)
            
            # 生成段落摘要
            part_summaries = []
            try:
                # 将文本分段处理
                paragraphs = [p.strip() for p in combined_text.split('\n\n') if p.strip()]
                for i, paragraph in enumerate(paragraphs[:5]):  # 只处理前5段
                    if len(paragraph) > 100:  # 只对较长的段落生成摘要
                        part_summary = self.llm_processor.generate_summary(paragraph)
                        part_summaries.append({
                            'paragraph_index': i + 1,
                            'original_length': len(paragraph),
                            'summary': part_summary
                        })
            except Exception as e:
                logger.warning(f"生成段落摘要失败: {e}")
            
            return {
                'summary': summary,
                'keywords': keywords,
                'hybrid_summary': hybrid_summary,
                'markdown_content': markdown_content,
                'part_summaries': part_summaries
            }
            
        except Exception as e:
            logger.error(f"生成摘要失败: {e}")
            return {
                'summary': '摘要生成失败',
                'keywords': [],
                'hybrid_summary': '分析失败',
                'markdown_content': '转换失败',
                'part_summaries': []
            }
    
    def _save_result(self, result: Dict, output_path: Path, output_name: str):
        """保存处理结果"""
        try:
            # 保存到pickle文件
            pickle_path = PICKLES_DIR / f"{output_name}_ppt.pkl"
            pickle_path.parent.mkdir(exist_ok=True)
            
            with open(pickle_path, 'wb') as f:
                pickle.dump(result, f)
            
            logger.info(f"结果已保存到: {pickle_path}")
            
            # 保存到JSON文件（可选）
            json_path = output_path / f"{output_name}_result.json"
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2, default=str)
            
            logger.info(f"结果已保存到: {json_path}")
            
        except Exception as e:
            logger.error(f"保存结果失败: {e}")
    
    def _cleanup_temp_files(self, output_path: Path):
        """清理临时文件"""
        try:
            # 可以在这里添加临时文件清理逻辑
            pass
        except Exception as e:
            logger.error(f"清理临时文件失败: {e}")
    
    def get_processing_status(self) -> Dict:
        """获取处理状态"""
        return {
            'supported_formats': self.supported_formats,
            'pptx_available': PPTX_AVAILABLE,
            'ppt_available': PPT_AVAILABLE,
            'llm_available': hasattr(self.llm_processor, 'llm') and self.llm_processor.llm is not None
        }
